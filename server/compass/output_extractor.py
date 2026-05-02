"""Extract text bodies from saved-output files for Compass audits.

The auto-audit watcher uses this to fold actual document content into
the artifact passed to `audit_work`, instead of auditing only path +
size. Output files are infrequent but high-stakes (binary deliverables
the human consumes), so spending the LLM tokens to actually read the
body is the right tradeoff (compass-specs §5.5, Tier B).

Extraction discipline:

  - **Lazy imports.** Each office-format parser (`pypdf`, `python-docx`,
    `openpyxl`, `python-pptx`) is imported inside the extractor that
    needs it. Missing libs degrade gracefully — the watcher falls back
    to path-only audit instead of crashing.
  - **Per-extractor exception isolation.** A malformed PDF must not
    take down a watcher iteration. `extract_body` returns `None` on
    any failure; the caller treats that as "no body, audit on metadata
    only".
  - **Output cap.** Bodies longer than `MAX_BODY_CHARS` get truncated
    with a marker. Audit verdicts mostly hinge on the doc's framing /
    main claims — the head usually carries enough signal, and tails
    blow up token costs without proportional value.
  - **Format families:**
      - **Text-native** (`md` / `txt` / `csv` / `tsv` / `html` / `json`):
        read the file as UTF-8.
      - **Office** (`pdf` / `docx` / `xlsx` / `pptx`): parse via the
        matching pure-Python lib.
      - **Archives** (`zip` / `tar` / `gz`): list filenames, no recursion.
        Often as informative as the contents would be (`payroll_q1.zip`
        with `salary_2026.csv` inside is a strong signal).
      - **Images** (`png` / `jpg` / etc.): skipped — Tier C (vision)
        is deferred. Caller falls back to path-only.
      - **Unknown**: same fallback as images.

The extractor never reads files outside the configured outputs dir.
Callers pass an already-resolved local path (`OUTPUTS_DIR / relpath`)
and `coord_save_output` rejects `..` segments at write time, but
`extract_body` enforces the boundary one more time as defense-in-depth:
a path that doesn't resolve under `OUTPUTS_DIR` returns `None` without
opening the file.
"""

from __future__ import annotations

import gzip
import logging
import sys
import tarfile
import zipfile
from pathlib import Path

logger = logging.getLogger("harness.compass.output_extractor")
if not logger.handlers:
    h = logging.StreamHandler(sys.stdout)
    h.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(name)s | %(message)s"))
    logger.addHandler(h)
    logger.setLevel(logging.INFO)


# Maximum characters of body text included in the audit artifact. Bigger
# documents are truncated with a marker. ~16k chars ≈ ~4k tokens, which
# leaves plenty of context-window headroom for the lattice + truth + the
# audit verdict response.
MAX_BODY_CHARS = 16_000

# Format families.
TEXT_NATIVE = {".md", ".markdown", ".txt", ".csv", ".tsv", ".html", ".htm", ".json"}
ARCHIVE = {".zip", ".tar", ".gz"}
# Skipped (Tier C / vision): {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}
# Unsupported: anything else falls through to path-only audit.


def extract_body(path: Path) -> str | None:
    """Return extracted text body for `path`, truncated to `MAX_BODY_CHARS`.
    Returns `None` when the format isn't body-extractable (image, unknown
    extension, missing parser dep) OR when extraction failed OR when the
    path resolves outside the configured outputs dir — in all three cases
    the caller falls back to path-only audit. Returns an empty string
    only when the file genuinely has no extractable text (e.g. a
    zero-byte CSV)."""
    if not path.is_file():
        return None
    if not _under_outputs_dir(path):
        logger.warning(
            "compass.output_extractor: refusing to read path outside outputs dir: %s",
            path,
        )
        return None
    ext = path.suffix.lower()
    try:
        if ext in TEXT_NATIVE:
            return _truncate(_extract_text(path))
        if ext == ".pdf":
            return _truncate_or_none(_extract_pdf(path))
        if ext == ".docx":
            return _truncate_or_none(_extract_docx(path))
        if ext == ".xlsx":
            return _truncate_or_none(_extract_xlsx(path))
        if ext == ".pptx":
            return _truncate_or_none(_extract_pptx(path))
        if ext in ARCHIVE:
            return _truncate(_extract_archive(path))
    except Exception:
        logger.exception(
            "compass.output_extractor: %s extraction failed for %s",
            ext, path.name,
        )
        return None
    # Image or unknown — fall back to path-only audit.
    return None


# ---------------------------------------------------------------- helpers


def _under_outputs_dir(path: Path) -> bool:
    """Return True iff `path` resolves to a location inside the
    configured `OUTPUTS_DIR`. Defense-in-depth boundary check —
    `coord_save_output` already rejects `..` segments at write time,
    so a real user can't bypass this through normal flow; the check
    catches accidental misuse from inside the harness (e.g. a future
    refactor that wires the extractor into a different caller).

    Lazy import of `outputs.OUTPUTS_DIR` so this module doesn't
    drag the outputs module into every test that touches the
    extractor.
    """
    from server.outputs import OUTPUTS_DIR  # noqa: PLC0415

    try:
        resolved = path.resolve()
        outputs_root = OUTPUTS_DIR.resolve()
    except OSError:
        return False
    try:
        resolved.relative_to(outputs_root)
    except ValueError:
        return False
    return True


def _truncate(s: str) -> str:
    s = (s or "").strip()
    if len(s) <= MAX_BODY_CHARS:
        return s
    return s[:MAX_BODY_CHARS] + f"\n\n[truncated — body is {len(s)} chars total]"


def _truncate_or_none(s: str | None) -> str | None:
    """Same as _truncate but preserves a None return — for parsers that
    can return None to signal a missing dep."""
    if s is None:
        return None
    return _truncate(s)


# ---------------------------------------------------------------- text-native


def _extract_text(path: Path) -> str:
    """Read a text-native file as UTF-8, replacing any decode errors so
    a single bad byte doesn't tank the audit."""
    return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------- PDF


def _extract_pdf(path: Path) -> str | None:
    """Extract concatenated page text via pypdf. Returns None if pypdf
    isn't installed (operator opted out of office-format auditing)."""
    try:
        from pypdf import PdfReader  # noqa: PLC0415
    except ImportError:
        logger.info("compass.output_extractor: pypdf not installed; PDFs audit by metadata only")
        return None
    reader = PdfReader(str(path))
    parts: list[str] = []
    char_budget = MAX_BODY_CHARS + 1000  # tiny over-read so truncation marker is honest
    for page in reader.pages:
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt:
            parts.append(txt)
        if sum(len(p) for p in parts) > char_budget:
            break
    return "\n\n".join(parts)


# ---------------------------------------------------------------- DOCX


def _extract_docx(path: Path) -> str | None:
    """Extract paragraph + table-cell text via python-docx."""
    try:
        from docx import Document  # noqa: PLC0415  (python-docx)
    except ImportError:
        logger.info("compass.output_extractor: python-docx not installed; .docx audit by metadata only")
        return None
    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_txt = " | ".join((cell.text or "").strip() for cell in row.cells)
            if row_txt.strip("| "):
                parts.append(row_txt)
    return "\n".join(parts)


# ---------------------------------------------------------------- XLSX


def _extract_xlsx(path: Path) -> str | None:
    """Extract sheet contents via openpyxl. Each sheet contributes a
    `## <SheetName>` header and TSV-style rows. read-only mode skips
    formatting parsing for speed.

    `wb.close()` runs in a finally block — read-only mode keeps a file
    handle open, and on Windows that handle stays attached to the file
    until close (or GC). A failure mid-iter shouldn't leak the FD into
    the watcher's process for an unbounded window."""
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    except ImportError:
        logger.info("compass.output_extractor: openpyxl not installed; .xlsx audit by metadata only")
        return None
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        char_budget = MAX_BODY_CHARS + 1000
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [
                    "" if v is None else str(v) for v in row
                ]
                if any(c.strip() for c in cells):
                    parts.append("\t".join(cells))
                if sum(len(p) for p in parts) > char_budget:
                    break
            if sum(len(p) for p in parts) > char_budget:
                break
        return "\n".join(parts)
    finally:
        wb.close()


# ---------------------------------------------------------------- PPTX


def _extract_pptx(path: Path) -> str | None:
    """Extract per-slide text frames via python-pptx. Each slide
    contributes a `### Slide N` header and the concatenated text of
    its shapes."""
    try:
        from pptx import Presentation  # noqa: PLC0415  (python-pptx)
    except ImportError:
        logger.info("compass.output_extractor: python-pptx not installed; .pptx audit by metadata only")
        return None
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"### Slide {i}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                txt = "".join(run.text or "" for run in para.runs).strip()
                if txt:
                    parts.append(txt)
    return "\n".join(parts)


# ---------------------------------------------------------------- archives


def _extract_archive(path: Path) -> str:
    """List filenames in a zip / tar / gz archive. Doesn't recurse into
    nested archives. Truncates the listing if the archive has thousands
    of entries — the head is enough to characterize the bundle.

    `.gz` files come in two flavors: tar-gzipped bundles
    (`report.tar.gz` / `bundle.tgz`) and plain gzipped single files
    (`report.csv.gz`). We try tarfile first since `.tar.gz` is the
    common archive shape; if that fails with a tar-format error we
    fall back to plain-gzip decompression and read the inner file as
    text. A real-world plain-gzip output (e.g. `data.csv.gz`) thus
    gets actual content into the audit instead of a path-only stub.

    On parser failure (corrupt archive, unreadable compression) lets
    the exception propagate so `extract_body`'s outer handler returns
    None (consistent with other extractors — corrupt = path-only audit
    fallback, not empty body)."""
    ext = path.suffix.lower()
    if ext == ".zip":
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
        return _format_archive_listing(path.name, names)
    if ext in (".tar", ".gz"):
        is_likely_tar = path.name.endswith((".tar", ".tar.gz", ".tgz"))
        if is_likely_tar:
            mode = "r:gz" if path.name.endswith((".tar.gz", ".tgz")) else "r"
            with tarfile.open(path, mode) as t:
                names = t.getnames()
            return _format_archive_listing(path.name, names)
        # Plain `.gz` (single gzipped file). Decompress + read as text
        # so the audit sees the inner content. Reads at most one body's
        # worth of bytes to bound memory.
        return _extract_plain_gzip(path)
    return ""


def _format_archive_listing(name: str, names: list[str]) -> str:
    parts: list[str] = [f"[archive listing — {name}]", f"({len(names)} entries)"]
    max_entries = 200
    for entry in names[:max_entries]:
        parts.append(f"  {entry}")
    if len(names) > max_entries:
        parts.append(f"  … and {len(names) - max_entries} more")
    return "\n".join(parts)


def _extract_plain_gzip(path: Path) -> str:
    """Decompress a single-file `.gz` and return the inner text body
    (UTF-8 with replacement). Bounded read — at most a body's worth
    plus a tiny over-read so the truncation marker downstream is
    honest. Returns a `[gzip body — name]` header so the LLM knows
    this content came from a compressed file rather than the original
    extension."""
    inner_name = path.stem  # `report.csv.gz` → `report.csv`
    char_budget = MAX_BODY_CHARS + 1000
    with gzip.open(path, "rb") as gz:
        raw = gz.read(char_budget)
    text = raw.decode("utf-8", errors="replace")
    return f"[gzip body — {inner_name}]\n{text}"


__all__ = [
    "MAX_BODY_CHARS",
    "TEXT_NATIVE",
    "ARCHIVE",
    "extract_body",
]
