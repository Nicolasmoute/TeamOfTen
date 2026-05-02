"""Tests for `server.compass.output_extractor`.

Format-specific text extractors for the Compass auto-audit watcher
(Tier B, compass-specs §5.5). Each extractor must:

  - Return extracted body text for known formats with content.
  - Return None for image / unknown formats (path-only fallback).
  - Truncate output longer than `MAX_BODY_CHARS` with a marker.
  - Survive parser failures without raising — return None instead.
"""

from __future__ import annotations

import gzip
import io  # noqa: F401  (kept for tests that build in-memory archives)
import zipfile
from pathlib import Path

import pytest

from server.compass import output_extractor as oe


# Each test writes its file under `tmp_path`. The extractor refuses to
# read paths outside `OUTPUTS_DIR` (defense-in-depth boundary check), so
# we point OUTPUTS_DIR at the per-test tmp_path before each test runs.
# Without this fixture every test would hit the boundary and return None
# regardless of format.
@pytest.fixture(autouse=True)
def _outputs_dir_to_tmp(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    from server import outputs as outmod

    monkeypatch.setattr(outmod, "OUTPUTS_DIR", tmp_path)


# ---------------------------------------------------------------- text-native


def test_extract_md_returns_body(tmp_path: Path) -> None:
    f = tmp_path / "report.md"
    f.write_text("# Report\n\nKey finding: per-task billing wins.", encoding="utf-8")
    body = oe.extract_body(f)
    assert body is not None
    assert "Key finding" in body
    assert "per-task billing" in body


def test_extract_txt_returns_body(tmp_path: Path) -> None:
    f = tmp_path / "notes.txt"
    f.write_text("plain text content", encoding="utf-8")
    body = oe.extract_body(f)
    assert body == "plain text content"


def test_extract_csv_returns_body(tmp_path: Path) -> None:
    f = tmp_path / "data.csv"
    f.write_text("name,price\nfoo,10\nbar,20\n", encoding="utf-8")
    body = oe.extract_body(f)
    assert body is not None
    assert "name,price" in body


def test_extract_html_returns_body(tmp_path: Path) -> None:
    """HTML is treated as text-native — we don't strip tags. Compass
    can reason about the content with the markup in place."""
    f = tmp_path / "page.html"
    f.write_text("<html><body><h1>Pricing</h1></body></html>", encoding="utf-8")
    body = oe.extract_body(f)
    assert body is not None
    assert "Pricing" in body


def test_extract_json_returns_body(tmp_path: Path) -> None:
    f = tmp_path / "config.json"
    f.write_text('{"plan": "pro", "seats": 5}', encoding="utf-8")
    body = oe.extract_body(f)
    assert body is not None
    assert '"plan"' in body


def test_extract_text_truncates_long_content(tmp_path: Path) -> None:
    f = tmp_path / "huge.txt"
    f.write_text("X" * (oe.MAX_BODY_CHARS + 5000), encoding="utf-8")
    body = oe.extract_body(f)
    assert body is not None
    assert "[truncated" in body
    assert len(body) < oe.MAX_BODY_CHARS + 200  # truncation kept it bounded


def test_text_decode_errors_are_replaced(tmp_path: Path) -> None:
    """A single bad byte must not tank the audit — we read with
    `errors='replace'`."""
    f = tmp_path / "broken.txt"
    f.write_bytes(b"valid \xff\xfe still readable")
    body = oe.extract_body(f)
    assert body is not None
    assert "valid" in body
    assert "still readable" in body


# ---------------------------------------------------------------- PDF


def test_extract_pdf_returns_body(tmp_path: Path) -> None:
    pypdf = pytest.importorskip("pypdf")
    from pypdf import PdfWriter

    f = tmp_path / "report.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with open(f, "wb") as out:
        writer.write(out)

    body = oe.extract_body(f)
    # Empty page, but extraction must succeed (returns "" or whitespace)
    # — the important thing is no crash and body is not None.
    assert body is not None


def test_extract_pdf_missing_dep_returns_none(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """If pypdf isn't installed, audit degrades to path-only."""
    f = tmp_path / "doc.pdf"
    f.write_bytes(b"%PDF-1.4\n")  # tiny stub, parser would fail anyway

    import builtins

    real_import = builtins.__import__

    def _no_pypdf(name: str, *args, **kwargs):
        if name == "pypdf":
            raise ImportError("simulated missing pypdf")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", _no_pypdf)
    body = oe.extract_body(f)
    assert body is None


def test_extract_corrupt_pdf_returns_none(tmp_path: Path) -> None:
    f = tmp_path / "broken.pdf"
    f.write_bytes(b"not really a pdf at all")
    body = oe.extract_body(f)
    assert body is None  # extractor caught the error and returned None


# ---------------------------------------------------------------- DOCX


def test_extract_docx_returns_body(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    from docx import Document

    f = tmp_path / "report.docx"
    doc = Document()
    doc.add_heading("Pricing analysis", level=1)
    doc.add_paragraph("Per-task billing is the recommended model.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "tier"
    table.cell(0, 1).text = "price"
    table.cell(1, 0).text = "pro"
    table.cell(1, 1).text = "$50"
    doc.save(str(f))

    body = oe.extract_body(f)
    assert body is not None
    assert "Pricing analysis" in body
    assert "Per-task billing" in body
    assert "pro | $50" in body


def test_extract_docx_missing_dep_returns_none(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    f = tmp_path / "doc.docx"
    f.write_bytes(b"PK")  # zip header stub

    import builtins

    real_import = builtins.__import__

    def _no_docx(name: str, *args, **kwargs):
        if name == "docx":
            raise ImportError("simulated missing python-docx")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", _no_docx)
    body = oe.extract_body(f)
    assert body is None


# ---------------------------------------------------------------- XLSX


def test_extract_xlsx_returns_body(tmp_path: Path) -> None:
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    f = tmp_path / "metrics.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Q1"
    ws.append(["metric", "value"])
    ws.append(["revenue", 12345])
    ws.append(["customers", 50])
    wb.save(str(f))

    body = oe.extract_body(f)
    assert body is not None
    assert "## Q1" in body
    assert "metric\tvalue" in body
    assert "revenue\t12345" in body


# ---------------------------------------------------------------- PPTX


def test_extract_pptx_returns_body(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    f = tmp_path / "deck.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Brand strategy"
    body_shape = slide.placeholders[1]
    body_shape.text = "Voice: plain and technical, not warm-conversational."
    prs.save(str(f))

    body = oe.extract_body(f)
    assert body is not None
    assert "### Slide 1" in body
    assert "Brand strategy" in body
    assert "plain and technical" in body


# ---------------------------------------------------------------- archives


def test_extract_zip_lists_filenames(tmp_path: Path) -> None:
    f = tmp_path / "bundle.zip"
    with zipfile.ZipFile(f, "w") as z:
        z.writestr("salary_2026.csv", "name,salary\nalice,100000\n")
        z.writestr("readme.md", "# bundle")
    body = oe.extract_body(f)
    assert body is not None
    assert "[archive listing — bundle.zip]" in body
    assert "salary_2026.csv" in body
    assert "readme.md" in body
    assert "(2 entries)" in body


def test_extract_corrupt_zip_returns_none(tmp_path: Path) -> None:
    f = tmp_path / "broken.zip"
    f.write_bytes(b"not really a zip")
    body = oe.extract_body(f)
    assert body is None


def test_archive_listing_truncates_huge_archive(tmp_path: Path) -> None:
    f = tmp_path / "many.zip"
    with zipfile.ZipFile(f, "w") as z:
        for i in range(250):  # exceeds the 200 max_entries cap
            z.writestr(f"file_{i:04d}.txt", "x")
    body = oe.extract_body(f)
    assert body is not None
    assert "(250 entries)" in body
    assert "and 50 more" in body


# ---------------------------------------------------------------- fallback


def test_image_extension_returns_none(tmp_path: Path) -> None:
    """PNG / JPG / etc. fall through to path-only audit (Tier C / vision
    is deferred)."""
    f = tmp_path / "chart.png"
    f.write_bytes(b"\x89PNG\r\n\x1a\n")
    assert oe.extract_body(f) is None


def test_unknown_extension_returns_none(tmp_path: Path) -> None:
    f = tmp_path / "weird.xyz"
    f.write_text("some data")
    assert oe.extract_body(f) is None


def test_missing_file_returns_none(tmp_path: Path) -> None:
    """A path that doesn't exist on disk shouldn't raise — just signals
    'no body, audit on metadata only'."""
    f = tmp_path / "ghost.md"
    assert oe.extract_body(f) is None


def test_directory_path_returns_none(tmp_path: Path) -> None:
    """Passing a directory path (not a file) is a defensive case —
    returns None instead of trying to read it."""
    d = tmp_path / "subdir"
    d.mkdir()
    assert oe.extract_body(d) is None


# ---------------------------------------------------------------- OUTPUTS_DIR boundary


def test_path_outside_outputs_dir_returns_none(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Defense-in-depth: a path that doesn't resolve under `OUTPUTS_DIR`
    must not be read, even if the file exists. The watcher's normal
    flow can't hit this (paths are composed as `OUTPUTS_DIR / relpath`
    and `coord_save_output` rejects `..`), but the check guards against
    a future caller misusing the extractor."""
    from server import outputs as outmod

    # Point OUTPUTS_DIR at a sibling directory so `tmp_path / "report.md"`
    # falls outside it. Override the autouse fixture's patch.
    other = tmp_path / "elsewhere"
    other.mkdir()
    monkeypatch.setattr(outmod, "OUTPUTS_DIR", other)

    f = tmp_path / "report.md"
    f.write_text("real content the extractor would normally happily read")
    assert oe.extract_body(f) is None


def test_path_inside_outputs_dir_subdir_works(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """Subdirectories under OUTPUTS_DIR are still readable — the
    boundary check uses `relative_to`, not strict-equal-parent."""
    from server import outputs as outmod

    monkeypatch.setattr(outmod, "OUTPUTS_DIR", tmp_path)
    sub = tmp_path / "reports" / "2026"
    sub.mkdir(parents=True)
    f = sub / "q1.md"
    f.write_text("Q1 results here")
    body = oe.extract_body(f)
    assert body is not None
    assert "Q1 results here" in body


# ---------------------------------------------------------------- plain gzip


def test_plain_gzip_extracts_inner_text(tmp_path: Path) -> None:
    """A bare `.gz` file (single gzipped file, NOT a tarball) gets
    decompressed and the inner content is included in the audit. Was
    previously a path-only fallback — now real content lands."""
    f = tmp_path / "report.csv.gz"
    inner = "name,price\nfoo,10\nbar,20\n"
    with gzip.open(f, "wb") as gz:
        gz.write(inner.encode("utf-8"))
    body = oe.extract_body(f)
    assert body is not None
    assert "[gzip body — report.csv]" in body
    assert "name,price" in body
    assert "foo,10" in body


def test_tar_gz_still_lists_filenames(tmp_path: Path) -> None:
    """Regression: `.tar.gz` and `.tgz` continue to be treated as
    archives (filename listing), not as plain gzip — the dual-mode
    detection picks tar for those suffixes."""
    import tarfile

    f = tmp_path / "bundle.tar.gz"
    with tarfile.open(f, "w:gz") as t:
        inner = tmp_path / "_inner.txt"
        inner.write_text("hello")
        t.add(inner, arcname="hello.txt")
    body = oe.extract_body(f)
    assert body is not None
    assert "[archive listing — bundle.tar.gz]" in body
    assert "hello.txt" in body


def test_corrupt_gzip_returns_none(tmp_path: Path) -> None:
    """A `.gz` that's neither a valid tar nor a valid plain gzip falls
    back to None (path-only audit) instead of raising."""
    f = tmp_path / "broken.gz"
    f.write_bytes(b"not a gzip stream at all")
    assert oe.extract_body(f) is None


def test_plain_gzip_decode_errors_replaced(tmp_path: Path) -> None:
    """Non-UTF-8 bytes inside a plain gzip don't tank the audit — the
    inner text uses `errors='replace'` like the text-native extractor."""
    f = tmp_path / "binary-ish.gz"
    with gzip.open(f, "wb") as gz:
        gz.write(b"valid \xff\xfe still readable")
    body = oe.extract_body(f)
    assert body is not None
    assert "[gzip body —" in body
    assert "valid" in body
    assert "still readable" in body
